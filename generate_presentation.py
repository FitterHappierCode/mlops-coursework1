"""
Generate a Coursework 1 PowerPoint using images and captions in data/figures/

Creates:
 - CW1_presentation_[YourName]_[YourID].pptx  (placeholder name/ID - edit filename manually)
 - CW1_speaker_script_[YourName]_[YourID].txt
 - CW1_slide_notes_[YourName]_[YourID].md
 - SUBMISSION_CHECKLIST.md

This script embeds captions directly under each figure and writes detailed notes.
Replace [YourName] and [YourID] in filenames before final save, or edit the variables below.
"""

import os
import textwrap
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from datetime import timedelta

BASE = os.path.dirname(__file__)
FIG_DIR = os.path.join(BASE, "data", "figures")
FIG_MD = os.path.join(FIG_DIR, "FIGURES.md")

# Output filenames (placeholders)
NAME = "[YourName]"
STUDENT_ID = "[YourID]"
# If an existing PPTX is locked, create a new file with a timestamp suffix to avoid permission errors
from datetime import datetime as _dt
OUT_PPTX = f"CW1_presentation_{NAME}_{STUDENT_ID}_" + _dt.now().strftime("%Y%m%d%H%M%S") + ".pptx"
OUT_SCRIPT = f"CW1_speaker_script_{NAME}_{STUDENT_ID}.txt"
OUT_NOTES_MD = f"CW1_slide_notes_{NAME}_{STUDENT_ID}.md"
OUT_CHECK = "SUBMISSION_CHECKLIST.md"

# Presentation timing
TOTAL_SECONDS = 450  # 7.5 minutes

# Simple parser for FIGURES.md: map image filename -> caption (first paragraph after heading)
def parse_figures_md(md_path):
    captions = {}
    if not os.path.exists(md_path):
        return captions
    with open(md_path, "r", encoding="utf-8") as f:
        lines = f.read().splitlines()
    cur_img = None
    capture = []
    for i, line in enumerate(lines):
        if line.strip().startswith("## "):
            # flush previous
            if cur_img and capture:
                text = " ".join([l.strip() for l in capture]).strip()
                captions[cur_img] = text
            # new heading: expected format '## filename'
            cur_img = line.strip()[3:].strip()
            capture = []
        else:
            # collect paragraphs until next heading
            if cur_img is not None:
                if line.strip() == "---":
                    # treat separator as end
                    continue
                capture.append(line)
    # flush last
    if cur_img and capture:
        text = " ".join([l.strip() for l in capture]).strip()
        captions[cur_img] = text
    return captions

# Small helper to add a title slide
def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    tx = slide.placeholders[1]
    tx.text = subtitle
    return slide

# Add a slide with bullet points
def add_bullet_slide(prs, title, bullets, notes=None):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        p = body.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(18)
    if notes:
        add_notes(slide, notes)
    return slide

# Add notes text
def add_notes(slide, notes_text):
    notes = slide.notes_slide
    text_frame = notes.notes_text_frame
    text_frame.clear()
    text_frame.text = notes_text

# Add image + caption slide (image on left or top)
def add_image_caption_slide(prs, title, image_path, caption, notes=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # title
    left = Inches(0.4)
    top = Inches(0.2)
    width = Inches(12.0)
    tx = slide.shapes.add_textbox(left, top, width, Inches(0.6))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    # image
    img_left = Inches(0.6)
    img_top = Inches(1.0)
    img_width = Inches(8.5)
    img_height = Inches(4.5)
    if image_path and os.path.exists(image_path):
        try:
            pic = slide.shapes.add_picture(image_path, img_left, img_top, width=img_width)
        except Exception:
            t = slide.shapes.add_textbox(img_left, img_top, Inches(8), Inches(1))
            t.text_frame.text = f"(Image not found: {os.path.basename(image_path)})"
    else:
        # place code/text area instead
        code_box = slide.shapes.add_textbox(img_left, img_top, img_width, img_height)
        code_tf = code_box.text_frame
        code_tf.word_wrap = True
        code_tf.text = caption
        # since we used caption as the code, clear caption below
        caption = ''
    # caption box beneath image
    cap_left = img_left
    cap_top = img_top + img_height + Inches(0.1)
    cap_w = Inches(9.0)
    cap_h = Inches(1.0)
    caption_box = slide.shapes.add_textbox(cap_left, cap_top, cap_w, cap_h)
    cf = caption_box.text_frame
    cf.word_wrap = True
    t = cf.paragraphs[0]
    t.text = caption
    t.font.size = Pt(12)
    t.font.italic = True
    # small footer
    footer = slide.shapes.add_textbox(Inches(10.2), Inches(6.9), Inches(2.0), Inches(0.4))
    footer_tf = footer.text_frame
    footer_tf.text = "COM774 — Slide"
    footer_tf.paragraphs[0].font.size = Pt(10)
    if notes:
        add_notes(slide, notes)
    return slide

# Helper: extract short code snippets from pipeline.py
def extract_snippet(file_path, anchors, context_lines=1):
    """Return a short snippet from file_path containing the first anchor line and surrounding context."""
    if not os.path.exists(file_path):
        return ""
    with open(file_path, 'r', encoding='utf-8') as f:
        lines = f.read().splitlines()
    for i, line in enumerate(lines):
        for anchor in anchors:
            if anchor in line:
                start = max(0, i - context_lines)
                end = min(len(lines), i + context_lines + 1)
                return '\n'.join(lines[start:end])
    return ""


# Build presentation
def build_presentation():
    captions = parse_figures_md(FIG_MD)
    fig_files = [f for f in os.listdir(FIG_DIR) if f.lower().endswith('.png')]

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Title slide
    title = "Coursework 1: Dataset selection & preprocessing"
    subtitle = f"COM774 — Intelligence Engineering and Infrastructure\nPrepared: [Your Name] | Student ID: [Your ID]"
    add_title_slide(prs, title, subtitle)

    # Core slides
    add_bullet_slide(prs, "Introduction & Aim", [
        "Aim: Select and prepare a dataset for a DevOps-related ML problem.",
        "This coursework focuses on the Data Pipeline; CW2 will build the ML model."],
        notes="Introduce the coursework purpose, deliverables and link to CW2.")

    add_bullet_slide(prs, "Dataset Overview", [
        "Source: local repository 'data/incidents_clean.csv' (cleaned output).",
        "Schema: key columns include first_opened_at, final_priority, resolution_hours, assignment_group_mode."],
        notes="Show 1-2 sample rows and mention size and date range.")

    add_bullet_slide(prs, "Why this dataset?", [
        "Relevant to incident management and SLA monitoring in DevOps.",
        "Feasible for CW2 ML tasks (predict breach/priority/duration)."],
        notes="Explain practical value and link to real-world DevOps operations.")

    add_bullet_slide(prs, "Data Quality Issues Found", [
        "Missing values in priority and timestamps; inconsistent assignment_group labels.",
        "Outliers in resolution hours and negative/invalid timestamps."],
        notes="Quantify counts where possible (use summary outputs).")

    add_bullet_slide(prs, "Cleaning & Preprocessing Steps", [
        "Drop exact duplicates and rows without opened timestamp.",
        "Normalize priority labels; parse timestamps; compute resolution_hours from timestamps.",
        "Clip extreme resolution outliers, encode categories, create SLA flags."],
        notes="Reference pipeline.py and explain rationale for each step.")

    # Add code snippet slides (short, annotated)
    pipeline_path = os.path.join(BASE, 'pipeline.py')
    # snippet 1: parse_dt and date parsing
    snip1 = extract_snippet(pipeline_path, ['def parse_dt', 'pd.to_datetime'])
    note1 = 'Parse mixed-format date/time strings and coerce failures to NaT to make downstream duration calculations safe.'
    if snip1:
        add_image_caption_slide(prs, 'Code: parse timestamps (snippet)', '', snip1 + '\n\n' + note1, notes=note1)

    # snippet 2: compute resolution_hours
    snip2 = extract_snippet(pipeline_path, ['resolution_hours_from_time', 'dt.total_seconds'], context_lines=2)
    note2 = 'Compute duration in hours from timestamps; prefer computed duration when available and coerce original values to numeric.'
    if snip2:
        add_image_caption_slide(prs, 'Code: compute resolution_hours (snippet)', '', snip2 + '\n\n' + note2, notes=note2)

    # snippet 3: normalize priority mapping
    snip3 = extract_snippet(pipeline_path, ['priority_map', 'def normalise_priority'], context_lines=3)
    note3 = 'Map many variants of priority labels to canonical categories to ensure consistent modelling.'
    if snip3:
        add_image_caption_slide(prs, 'Code: normalise priority labels (snippet)', '', snip3 + '\n\n' + note3, notes=note3)

    # Visual results slides (reuse earlier approach)
    vis_order = [
        'incidents_over_time.png',
        'missingness_heatmap.png',
        'incidents_per_priority.png',
        'sla_rate_by_month.png',
        'resolution_by_priority.png',
        'incidents_dirty_vs_clean.png'
    ]
    for i in range(0, len(vis_order), 3):
        group = vis_order[i:i+3]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.6)).text_frame.text = 'Visual results'
        y = 1.0
        for img in group:
            img_path = os.path.join(FIG_DIR, img)
            cap = captions.get(img, '')
            try:
                slide.shapes.add_picture(img_path, Inches(0.5), Inches(y), width=Inches(4.0))
            except Exception:
                slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(4.0), Inches(1.0)).text_frame.text = f'(Missing: {img})'
            cf = slide.shapes.add_textbox(Inches(0.5), Inches(y+2.2), Inches(4.0), Inches(0.8)).text_frame
            cf.word_wrap = True
            cf.text = cap
            y += 2.8
        add_notes(slide, 'Discuss interpretation and implication of each figure; refer to captions under images.')

    # Limitations, Next steps, Conclusion, Integrity, Submission checklist
    add_bullet_slide(prs, 'Limitations & Assumptions', [
        'Data biases, gaps in time range, synthetic anomalies if used.',
        'Decisions (clipping, encoding) may affect model behaviour in CW2.'], notes='Be candid about limitations and how to mitigate.')

    add_bullet_slide(prs, 'Next steps for CW2', [
        'Define target (SLA breach prediction) and candidate features.',
        'Model evaluation plan: metrics, baseline, retraining triggers.'], notes='Outline model approach and MLOps considerations.')

    add_bullet_slide(prs, 'Conclusion & Key Takeaways', [
        'Dataset is feasible and relevant to incident management.',
        'Cleaning steps improved quality and produced key features for ML.'], notes='Summarise and provide call-to-action for CW2.')

    integrity = (
        'I declare that this is all my own work. Any material I have referred to has been accurately referenced and any contribution of Artificial Intelligence technology has been fully acknowledged.'
    )
    add_bullet_slide(prs, 'Academic Integrity Declaration', [integrity], notes='Include full coursework declaration from brief.')

    add_bullet_slide(prs, 'Submission Checklist', [
        'File: CW1_presentation_[YourName]_[YourID].pptx (ensure name & ID on first slide).',
        'Embedded video and audio checked; length ≈7.5 minutes.',
        'References and academic integrity slide present.'], notes='Final checks before upload.')

    add_bullet_slide(prs, 'Reproducibility & Commands', [
        'Run: python pipeline.py -> data/incidents_clean.csv',
        'Run: python make_more_visuals.py -> data/figures/*.png',
        'Python packages: see requirements or virtualenv used.'], notes='Include exact versions in README if required.')

    # Appendix: full pipeline snippet
    if os.path.exists(pipeline_path):
        with open(pipeline_path, 'r', encoding='utf-8') as pf:
            full = '\n'.join(pf.read().splitlines()[:120])  # first 120 lines to keep length reasonable
        add_image_caption_slide(prs, 'Appendix: pipeline.py (start)', '', full, notes='Full pipeline available in repository (pipeline.py)')

    out_path = os.path.join(BASE, OUT_PPTX)
    prs.save(out_path)

    # Build speaker script / slide notes text
    all_slides = prs.slides
    n = len(all_slides)
    per_slide = TOTAL_SECONDS // n
    extra = TOTAL_SECONDS - per_slide * n
    t = 0
    with open(os.path.join(BASE, OUT_SCRIPT), 'w', encoding='utf-8') as sfile, open(os.path.join(BASE, OUT_NOTES_MD), 'w', encoding='utf-8') as nfile:
        sfile.write(f"Timed speaker script for {OUT_PPTX}\nTotal time: {TOTAL_SECONDS} seconds (~7.5 minutes)\n\n")
        for idx, slide in enumerate(all_slides, start=1):
            duration = per_slide + (1 if idx <= extra else 0)
            start = timedelta(seconds=t)
            end = timedelta(seconds=t+duration)
            notes_text = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
            title_text = slide.shapes.title.text if slide.shapes.title else f"Slide {idx}"
            sfile.write(f"[{str(start)} - {str(end)}] Slide {idx}: {title_text}\n")
            sfile.write(notes_text + "\n\n")
            nfile.write(f"## Slide {idx}: {title_text}\n\n")
            nfile.write(notes_text + "\n\n")
            t += duration

    # Write submission checklist file
    with open(os.path.join(BASE, OUT_CHECK), 'w', encoding='utf-8') as cf:
        cf.write("SUBMISSION CHECKLIST\n\n")
        cf.write("1. Ensure first slide contains your name and student ID.\n")
        cf.write("2. Confirm file saved as CW1_presentation_[YourName]_[YourID].pptx\n")
        cf.write("3. Check embedded figures and captions appear beneath images.\n")
        cf.write("4. Verify total spoken length is ~7.5 minutes using the script.\n")
        cf.write("5. Upload to Blackboard Assessment 1 before the deadline.\n")

    print("Created presentation and supporting files:")
    print(" -", out_path)
    print(" -", os.path.join(BASE, OUT_SCRIPT))
    print(" -", os.path.join(BASE, OUT_NOTES_MD))
    print(" -", os.path.join(BASE, OUT_CHECK))


if __name__ == '__main__':
    build_presentation()
